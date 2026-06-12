"""Conversion utilities: PDF ↔ Word/Excel/PPT/Images"""
import fitz  # PyMuPDF
import pdfplumber
import os
import shutil
from PIL import Image
import pandas as pd
from pptx import Presentation
from pptx.util import Inches


# ---------------------------------------------------------------------------
# PDF → X
# ---------------------------------------------------------------------------

def pdf_to_images(input_path: str, output_dir: str, fmt: str = "png", dpi: int = 150) -> list[str]:
    """Convert each PDF page to an image."""
    os.makedirs(output_dir, exist_ok=True)
    results = []
    base = os.path.splitext(os.path.basename(input_path))[0]
    with fitz.open(input_path) as doc:
        for i, page in enumerate(doc):
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            out_path = os.path.join(output_dir, f"{base}_page{i+1}.{fmt}")
            pix.save(out_path)
            results.append(out_path)
    return results


def pdf_to_word(input_path: str, output_path: str) -> str:
    """Convert PDF to DOCX using pdf2docx."""
    from pdf2docx import Converter
    cv = Converter(input_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
    return output_path


def pdf_to_excel(input_path: str, output_path: str) -> str:
    """Extract tables from PDF and save to Excel using pdfplumber."""
    all_tables = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            for table in page.extract_tables():
                df = pd.DataFrame(table[1:], columns=table[0] if table else None)
                all_tables.append(df)
    if not all_tables:
        raise ValueError("No tables found in the PDF.")
    with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
        for idx, df in enumerate(all_tables):
            df.to_excel(writer, sheet_name=f"Table_{idx+1}", index=False)
    return output_path


def pdf_to_ppt(input_path: str, output_path: str, dpi: int = 150) -> str:
    """Convert PDF pages to PPT slides (each page as an image)."""
    tmp_dir = output_path + "_tmp_imgs"
    images = pdf_to_images(input_path, tmp_dir, fmt="png", dpi=dpi)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for img_path in images:
        slide = prs.slides.add_slide(blank_layout)
        with Image.open(img_path) as im:
            w, h = im.size
        aspect = h / w
        slide_w = prs.slide_width
        slide_h = int(slide_w * aspect)
        prs.slide_height = slide_h
        slide.shapes.add_picture(img_path, 0, 0, slide_w, slide_h)
    prs.save(output_path)
    shutil.rmtree(tmp_dir, ignore_errors=True)
    return output_path


# ---------------------------------------------------------------------------
# X → PDF
# ---------------------------------------------------------------------------

def word_to_pdf(input_path: str, output_path: str) -> str:
    """Convert DOCX to PDF using docx2pdf (requires MS Word on Windows)."""
    from docx2pdf import convert
    convert(input_path, output_path)
    return output_path


def excel_to_pdf(input_path: str, output_path: str) -> str:
    """Convert XLSX to PDF using openpyxl + reportlab."""
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
    from reportlab.lib import colors
    import openpyxl

    wb = openpyxl.load_workbook(input_path, data_only=True)
    doc = SimpleDocTemplate(output_path, pagesize=landscape(A4))
    story = []

    for sheet in wb.worksheets:
        data = [[str(cell.value) if cell.value is not None else "" for cell in row]
                for row in sheet.iter_rows()]
        if not data:
            continue
        t = Table(data, repeatRows=1)
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#7c3aed")),
            ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",   (0, 0), (-1, -1), 8),
            ("GRID",       (0, 0), (-1, -1), 0.5, colors.grey),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f3f0ff")]),
        ]))
        story.append(t)

    doc.build(story)
    return output_path


def ppt_to_pdf(input_path: str, output_path: str) -> str:
    """Convert PPTX to PDF using docx2pdf (requires MS Office on Windows)."""
    from docx2pdf import convert
    convert(input_path, output_path)
    return output_path


def images_to_pdf(input_paths: list[str], output_path: str) -> str:
    """Combine multiple images into a single PDF using PyMuPDF."""
    doc = fitz.open()
    for img_path in input_paths:
        with Image.open(img_path) as im:
            rgb = im.convert("RGB")
            tmp = img_path + "_tmp.jpg"
            rgb.save(tmp, "JPEG")
        img_doc = fitz.open(tmp)
        pdf_bytes = img_doc.convert_to_pdf()
        img_doc.close()
        os.remove(tmp)
        img_pdf = fitz.open("pdf", pdf_bytes)
        doc.insert_pdf(img_pdf)
    doc.save(output_path)
    doc.close()
    return output_path
